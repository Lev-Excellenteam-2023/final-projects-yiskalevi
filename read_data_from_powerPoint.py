import os
from pptx import Presentation

# def openFile(path: str):
#     """
#     Gets a path to open a file and opens the file, checks that it is correct and exists.
#     """
#     if os.path.isfile(path) and path.endswith('.pptx'):
#         return True
#     else:
#         raise ValueError("Invalid file path or file format. Please provide a valid .pptx file.")

def read_slides_from_presentation(file) -> list:
    """
    Reads the information from the presentation and returns a list where each element represents a slide.
    """
    presentation = Presentation(file)
    slides = []
    for slide in presentation.slides:
        text_boxes = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                text = ""
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        text += run.text
                text_boxes.append(text.strip())
        slides.append(text_boxes)
    return slides

# def read_data_from_slides(lst: list) -> str:
#     """
#     Extracts all the text information from the slides.
#     """
#     text = ""
#     for slide in lst:
#         for box in slide:
#             text += box + "\n"
#     return text.strip()
