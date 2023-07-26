from pptx import Presentation


def read_slides_from_presentation(file) -> list:
    """
    Reads the information from the presentation and returns a list where each element represents a slide.
    """
    presentation = Presentation(file)
    slides = []
    for slide in presentation.slides:
        text_boxes = []
        text_in_file=""
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                text = ""
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        text += run.text
                text_boxes.append(text.strip())
                text_in_file+=text.strip()
        slides.append(text_in_file)
    return slides

